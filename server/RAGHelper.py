import hashlib
import os
import glob
from tqdm import tqdm
import json
import jq
from pptx import Presentation

import tiktoken

from LLMHelper import LLMHelper

from sentence_transformers import SentenceTransformer
from PostgresHybridRetriever import PostgresHybridRetriever

from docling.document_converter import DocumentConverter

from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_experimental.text_splitter import SemanticChunker
from ParagraphChunker import ParagraphChunker

from Reranker import Reranker

from provenance import compute_llm_provenance, compute_rerank_provenance, DocumentSimilarityAttribution

class RAGHelper:
    """
    A helper class to manage retrieval-augmented generation (RAG) processes,
    including data loading, chunking, vector storage, and retrieval.
    """

    def __init__(self, logger, db_pool):
        """
        Initializes the RAGHelper class and loads environment variables.
        """
        self.logger = logger
        self.db_pool = db_pool

        # Set up docling
        self.converter = DocumentConverter()

        # Initialize the LLM and embeddings
        self.llm = LLMHelper(logger)
        self.embeddings = self.initialize_embeddings()

        # Set up the PostgresHybridRetriever
        self.retriever = PostgresHybridRetriever(self.db_pool)
        self.retriever.setup_database(self.embeddings.get_sentence_embedding_dimension())

        # Initialize the reranker
        if os.getenv("rerank") == "True":
            self.logger.info("Initializing reranker.")
            self.reranker = Reranker()

        # Load the data into the vector store
        self.splitter = self._initialize_text_splitter()
        if not self.retriever.has_data():
            self.load_data()
        
        # Provenance
        if os.getenv("provenance_method") == "similarity":
            self.similarity_attribution = DocumentSimilarityAttribution()
        
        # Summarization
        if os.getenv("use_summarization") == "True":
            self.tiktoken_encoder = tiktoken.encoding_for_model(os.getenv("summarization_encoder"))

    ############################
    ### Initialization functions
    ############################
    def reload_llm(self):
        """Reinitialise the LLM client (and reranker if enabled) after
        environment variables have been refreshed.  This is called from
        the /config PUT endpoint when ``reinitialize`` is requested."""
        self.logger.info("Reloading LLM client.")
        self.llm = LLMHelper(self.logger)

        # Reinitialise reranker if configured
        if os.getenv("rerank") == "True":
            self.logger.info("Reinitializing reranker.")
            self.reranker = Reranker()

        # Re-check provenance
        if os.getenv("provenance_method") == "similarity":
            self.similarity_attribution = DocumentSimilarityAttribution()

        # Summarization encoder
        if os.getenv("use_summarization") == "True":
            import tiktoken
            self.tiktoken_encoder = tiktoken.encoding_for_model(os.getenv("summarization_encoder"))

    def initialize_embeddings(self):
        """Initialize the embeddings based on the CPU/GPU configuration."""
        embedding_model = os.getenv("embedding_model")
        device = (
            "cpu"
            if os.getenv("embedding_cpu") == "True"
            else "cuda"
        )
        self.logger.info(f"Initializing embedding model {embedding_model} on device {device}.")
        return SentenceTransformer(embedding_model, device=device)
    
    def _initialize_text_splitter(self):
        """Initialize the text splitter based on the environment settings."""
        splitter_type = os.getenv("splitter")
        self.logger.info(f"Initializing {splitter_type} splitter.")
        if splitter_type == "RecursiveCharacterTextSplitter":
            return RecursiveCharacterTextSplitter(
            chunk_size=int(os.getenv("recursive_splitter_chunk_size")),
            chunk_overlap=int(os.getenv("recursive_splitter_chunk_overlap")),
            length_function=len,
            keep_separator=True,
            separators=[
                "\n \n",
                "\n\n",
                "\n",
                ".",
                "!",
                "?",
                " ",
                ",",
                "\u200b",
                "\uff0c",
                "\u3001",
                "\uff0e",
                "\u3002",
                "",
            ],
        )
        elif splitter_type == "SemanticChunker":
            return SemanticChunker(
                self.embeddings,
                breakpoint_threshold_type=os.getenv("semantic_chunker_breakpoint_threshold_type"),
                breakpoint_threshold_amount=os.getenv("semantic_chunker_breakpoint_threshold_amount"),
                number_of_chunks=os.getenv("semantic_chunker_number_of_chunks"),
            )
        elif splitter_type == "ParagraphChunker":
            return ParagraphChunker(
                max_chunk_size=int(os.getenv("paragraph_chunker_max_chunk_size")),
                paragraph_separator=os.getenv("paragraph_chunker_paragraph_separator")
            )
    
    ########################
    ### Data Loading Section
    ########################
    def load_data(self):
        """
        Loads data from various file types and chunks it into an ensemble retriever.
        """
        data_dir = os.getenv("data_directory")
        file_types = os.getenv("file_types").split(",")

        if "json" in file_types:
            jq_compiled = jq.compile(os.getenv("json_schema"))

        # Load all files in the data directory, recursively
        files = glob.glob(os.path.join(data_dir, "**"), recursive=True)
        documents = []
        with tqdm(
            total=len(files), desc="Reading in, chunking, and vectorizing documents"
        ) as pbar:
            for file in files:
                file_type = file.split(".")[-1]
                if os.path.isfile(file) and file_type in file_types:
                    # Load the document based on the file type
                    if file_type == "json":
                        with open(file, "r", encoding="utf-8") as f:
                            doc = json.load(f)
                            doc = jq_compiled.input(doc).first()
                            doc = json.dumps(doc)
                    elif file_type == "txt" or file_type == "xml":
                        with open(file, "r", encoding="utf-8") as f:
                            doc = f.read()
                    elif file_type == "csv":
                        import pandas as pd
                        df = pd.read_csv(file, encoding="utf-8", sep=os.getenv("csv_separator"))
                        json_data = df.to_dict(orient='records')
                        doc = json.dumps(json_data)
                    elif file_type == "pptx":
                        presentation = Presentation(file)
                        full_text = []
                        for slide in presentation.slides:
                            slide_text = []
                            for shape in slide.shapes:
                                if shape.has_text_frame:
                                    for paragraph in shape.text_frame.paragraphs:
                                        slide_text.append(paragraph.text)
                            full_text.append("\n".join(slide_text))
                        doc = "\n\n".join(full_text)
                    else:
                        doc = self.converter.convert(file).document.export_to_text()
                    
                    # Get the subfolder name of this document
                    subfolder = os.path.basename(os.path.dirname(file)).replace(os.getenv("data_directory"), "")
                    if not(file_type == "csv"):
                        # Chunk the document
                        chunks = self.splitter.split_text(doc)
                    else:
                        chunks = doc
                    
                    chunks = [{
                        "id": hashlib.md5(chunk.encode()).hexdigest(),
                        "embedding": self.embeddings.encode(chunk),
                        "content": chunk,
                        "metadata": json.dumps({
                            "source": file,
                            "dataset": subfolder
                        })
                    } for chunk in chunks]

                    # Insert the chunks into the vector store
                    documents.extend(chunks)
                pbar.update(1)
        
        self.logger.info(f"Writing {len(documents)} documents to the vector store.")
        documents = self._deduplicate_chunks(documents)
        self.retriever.add_documents(documents)
    
    def add_document(self, file_path, dataset):
        file_type = file_path.split(".")[-1]

        if file_type == "json":
            jq_compiled = jq.compile(os.getenv("json_schema"))

        # Load the document based on the file type
        documents = []
        file_types = os.getenv("file_types").split(",")
        if os.path.isfile(file_path) and file_type in file_types:
            # Load the document based on the file type
            if file_type == "json":
                with open(file_path, "r", encoding="utf-8") as f:
                    doc = json.load(f)
                    doc = jq_compiled.input(doc).first()
                    doc = json.dumps(doc)
            elif file_type == "txt" or file_type == "xml":
                with open(file_path, "r", encoding="utf-8") as f:
                    doc = f.read()
            elif file_type == "pptx":
                presentation = Presentation(file_path)
                full_text = []
                for slide in presentation.slides:
                    slide_text = []
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for paragraph in shape.text_frame.paragraphs:
                                slide_text.append(paragraph.text)
                    full_text.append("\n".join(slide_text))
                doc = "\n\n".join(full_text)
            else:
                doc = self.converter.convert(file_path).document.export_to_text()
            
            # Chunk the document
            chunks = self.splitter.split_text(doc)
            chunks = [{
                "id": hashlib.md5(chunk.encode()).hexdigest(),
                "embedding": self.embeddings.encode(chunk),
                "content": chunk,
                "metadata": json.dumps({
                    "source": file_path,
                    "dataset": dataset
                })
            } for chunk in chunks]

            # Insert the chunks into the vector store
            documents.extend(chunks)
        
        self.logger.info(f"Wrote document {file_path} to the vector store.")
        documents = self._deduplicate_chunks(documents)
        self.retriever.add_documents(documents)

    def _deduplicate_chunks(self, documents):
        return list({doc['id']: doc for doc in documents}.values())

    ##################
    ### Chat functions
    ##################
    def format_documents(self, docs):
        """
        Formats the documents for better readability.

        Args:
            docs (list): List of Document objects.

        Returns:
            str: Formatted string representation of documents.
        """
        doc_strings = []
        for i, doc in enumerate(docs):
            metadata_string = ", ".join(
                [f"{md}: {doc['metadata'][md]}" for md in doc['metadata'].keys()]
            )
            filename = doc['metadata']['source']
            doc_strings.append(
                f"[Document] *Filename* `{filename}`\n*Content*: {doc['content']}\n*Metadata* {metadata_string} [/Document]"
            )
        return "\n\n".join(doc_strings)

    def handle_documents(self, prompt, prompt_embedding, datasets, step_callback=None):
        # Reobtain documents with new question
        documents = self.retriever.get_relevant_documents(prompt, prompt_embedding, datasets)

        # Check if we need to apply the reranker and run it
        if os.getenv("rerank") == "True":
            if step_callback:
                step_callback(f"Reranking top {os.getenv('rerank_k')} documents...")
            self.logger.info("Reranking documents.")
            documents = self.reranker.rerank_documents(documents, prompt)[:int(os.getenv("rerank_k"))]
        else:
            documents = [{**document, "score": document['metadata']['distance']} for document in documents]
        
        return documents

    def compute_provenance_scores(self, prompt, documents, response):
        # Compute the provenance score
        provenance_scores = None
        if os.getenv("provenance_method") in ["rerank", "llm", "similarity"]:
            if os.getenv("provenance_method") == "rerank":
                provenance_scores = compute_rerank_provenance(self.reranker, prompt, documents, response)
            elif os.getenv("provenance_method") == "llm":
                provenance_scores = compute_llm_provenance(self.llm, prompt, documents, response)
            elif os.getenv("provenance_method") == "similarity":
                provenance_scores = self.similarity_attribution.compute_similarity(prompt, documents, response)

        return provenance_scores

    def handle_user_interaction(self, prompt, history, datasets):
        """
        Handle user interaction with the RAG system.
        """
        rewritten = None
        # Check if we need to fetch new documents
        fetch_new_documents = True
        if len(history) > 0:
            # Summarize the history if needed
            if os.getenv("use_summarization") == "True":
                # Convert the history to a string
                self.logger.info("Checking if we need to summarize the history.")
                history_string = "\n\n".join([f"{message['role']}: {message['content']}" for message in history])
                history_size = len(self.tiktoken_encoder.encode(history_string))
                if history_size > int(os.getenv("summarization_threshold")):
                    self.logger.info(f"Summarizing the history because it contains {history_size} tokens.")
                    (response, _) = self.llm.generate_response(
                        None,
                        os.getenv("summarization_query").format(history=history_string),
                        []
                    )
                    history = history[:1] + [{"role": "assistant", "content": response}]
            
            # Get the LLM response to see if we need to fetch new documents
            self.logger.info("History is not empty, checking if we need to fetch new documents.")
            (response, _) = self.llm.generate_response(
                None,
                os.getenv("rag_fetch_new_question").format(question=prompt),
                history
            )
            if response.lower().strip().startswith("no"):
                fetch_new_documents = False
        
        # Fetch new documents if needed
        documents = None
        if fetch_new_documents:
            # Apply hyde if needed
            if os.getenv("use_hyde") == "True":
                (response, _) = self.llm.generate_response(
                    None,
                    os.getenv("hyde_query").format(question=prompt),
                    []
                )
                prompt = response

            self.logger.info("Fetching new documents.")
            prompt_embedding = self.embeddings.encode(prompt)
            documents = self.handle_documents(prompt, prompt_embedding, datasets)

            # Check if the answer is in the documents or not
            if os.getenv("use_rewrite_loop") == "True" and not os.getenv("use_hyde") == "True":
                self.logger.info("Rewrite is enabled - checking if the fetched documents contain the answer.")
                (response, _) = self.llm.generate_response(
                    os.getenv("rewrite_query_instruction").format(context=self.format_documents(documents)),
                    os.getenv("rewrite_query_question").format(question=prompt),
                    []
                )
                if response.lower().strip().startswith("no"):
                    # Rewrite the query
                    self.logger.info("Rewrite is enabled and the answer is not in the documents - rewriting the query.")
                    (new_prompt, _) = self.llm.generate_response(
                        None,
                        os.getenv("rewrite_query_prompt").format(question=prompt, motivation=f"Can I find the answer in the documents: {response}"),
                        []
                    )
                    self.logger.info(f"Rewrite complete, original query: {prompt}, rewritten query: {new_prompt}")
                    rewritten = new_prompt
                    # Reobtain documents with new question
                    documents = self.handle_documents(new_prompt, prompt_embedding, datasets)
                else:
                    self.logger.info("Rewrite is enabled but the query is adequate.")
            else:
                self.logger.info("Rewrite is disabled - using the original query.")
        
        # Apply RE2 if turend on (but not in conjunction with hyde)
        if os.getenv("use_re2") == "True" and not os.getenv("use_hyde") == "True":
            prompt = f"{prompt}\n{os.getenv('re2_prompt')}\n{prompt}"
        
        provenance_scores = None

        # Get the LLM response
        if len(history) == 0:
            (response, new_history) = self.llm.generate_response(
                os.getenv("rag_instruction").format(context=self.format_documents(documents)),
                os.getenv("rag_question_initial").format(question=prompt),
                []
            )
            provenance_scores = self.compute_provenance_scores(prompt, documents, response)
        elif fetch_new_documents:
            # Add the documents to the system prompt and remove the previous system prompt
            (response, new_history) = self.llm.generate_response(
                os.getenv("rag_instruction").format(context=self.format_documents(documents)),
                os.getenv("rag_question_followup").format(question=prompt),
                [message for message in history if message["role"] != "system"]
            )
            provenance_scores = self.compute_provenance_scores(prompt, documents, response)
        else:
            # Keep the full history, with system prompt and previous documents
            (response, new_history) = self.llm.generate_response(
                None,
                os.getenv("rag_question_followup").format(question=prompt),
                history
            )
        
        # Add the response to the history
        new_history.append({"role": "assistant", "content": response})
        return (response, documents, fetch_new_documents, rewritten, new_history, provenance_scores)

    def handle_user_interaction_stream(self, prompt, history, datasets):
        """
        Streaming version of handle_user_interaction.
        Yields SSE-formatted events for each pipeline step and then streams the LLM response tokens.
        Event types:
          - ("step", step_name)   : pipeline step status
          - ("token", text)       : LLM response chunk
          - ("done", metadata)    : final metadata dict with history, documents, etc.
        """
        import json as _json

        rewritten = None
        fetch_new_documents = True

        # Summarization check
        if len(history) > 0:
            if os.getenv("use_summarization") == "True":
                yield ("step", "Checking if history needs summarization...")
                history_string = "\n\n".join([f"{message['role']}: {message['content']}" for message in history])
                history_size = len(self.tiktoken_encoder.encode(history_string))
                if history_size > int(os.getenv("summarization_threshold")):
                    yield ("step", "Summarizing conversation history...")
                    self.logger.info(f"Summarizing the history because it contains {history_size} tokens.")
                    (response, _) = self.llm.generate_response(
                        None,
                        os.getenv("summarization_query").format(history=history_string),
                        []
                    )
                    history = history[:1] + [{"role": "assistant", "content": response}]

            # Check if we need to fetch new documents
            yield ("step", "Checking if new documents are needed...")
            self.logger.info("History is not empty, checking if we need to fetch new documents.")
            (response, _) = self.llm.generate_response(
                None,
                os.getenv("rag_fetch_new_question").format(question=prompt),
                history
            )
            if response.lower().strip().startswith("no"):
                fetch_new_documents = False
                yield ("step", "Using existing context (no new retrieval needed).")

        # Fetch new documents if needed
        documents = None
        pending_steps = []  # collect steps from sub-calls that can't yield directly
        if fetch_new_documents:
            # HyDE
            if os.getenv("use_hyde") == "True":
                yield ("step", "Generating hypothetical document (HyDE)...")
                (response, _) = self.llm.generate_response(
                    None,
                    os.getenv("hyde_query").format(question=prompt),
                    []
                )
                prompt = response

            yield ("step", "Retrieving relevant documents...")
            self.logger.info("Fetching new documents.")
            prompt_embedding = self.embeddings.encode(prompt)
            documents = self.handle_documents(prompt, prompt_embedding, datasets, step_callback=lambda s: pending_steps.append(s))
            for s in pending_steps:
                yield ("step", s)
            pending_steps.clear()

            # Rewrite loop
            if os.getenv("use_rewrite_loop") == "True" and not os.getenv("use_hyde") == "True":
                yield ("step", "Checking if documents contain the answer...")
                self.logger.info("Rewrite is enabled - checking if the fetched documents contain the answer.")
                (response, _) = self.llm.generate_response(
                    os.getenv("rewrite_query_instruction").format(context=self.format_documents(documents)),
                    os.getenv("rewrite_query_question").format(question=prompt),
                    []
                )
                if response.lower().strip().startswith("no"):
                    yield ("step", "Rewriting query for better results...")
                    self.logger.info("Rewrite is enabled and the answer is not in the documents - rewriting the query.")
                    (new_prompt, _) = self.llm.generate_response(
                        None,
                        os.getenv("rewrite_query_prompt").format(question=prompt, motivation=f"Can I find the answer in the documents: {response}"),
                        []
                    )
                    self.logger.info(f"Rewrite complete, original query: {prompt}, rewritten query: {new_prompt}")
                    rewritten = new_prompt
                    yield ("step", "Re-retrieving documents with improved query...")
                    documents = self.handle_documents(new_prompt, prompt_embedding, datasets, step_callback=lambda s: pending_steps.append(s))
                    for s in pending_steps:
                        yield ("step", s)
                    pending_steps.clear()
                else:
                    yield ("step", "Documents look relevant, proceeding...")
                    self.logger.info("Rewrite is enabled but the query is adequate.")
            else:
                self.logger.info("Rewrite is disabled - using the original query.")

        # RE2
        if os.getenv("use_re2") == "True" and not os.getenv("use_hyde") == "True":
            yield ("step", "Applying RE2 (Re-reading) prompt enhancement...")
            prompt = f"{prompt}\n{os.getenv('re2_prompt')}\n{prompt}"

        # Send documents to the client before LLM generation
        if documents:
            yield ("documents", documents)

        # Stream the LLM response
        yield ("step", "Generating answer...")
        full_response_chunks = []

        if len(history) == 0:
            (stream, new_history) = self.llm.generate_response_stream(
                os.getenv("rag_instruction").format(context=self.format_documents(documents)),
                os.getenv("rag_question_initial").format(question=prompt),
                []
            )
        elif fetch_new_documents:
            (stream, new_history) = self.llm.generate_response_stream(
                os.getenv("rag_instruction").format(context=self.format_documents(documents)),
                os.getenv("rag_question_followup").format(question=prompt),
                [message for message in history if message["role"] != "system"]
            )
        else:
            (stream, new_history) = self.llm.generate_response_stream(
                None,
                os.getenv("rag_question_followup").format(question=prompt),
                history
            )

        for chunk in stream:
            full_response_chunks.append(chunk)
            yield ("token", chunk)

        response = "".join(full_response_chunks)

        # Compute provenance
        provenance_scores = None
        provenance_method = os.getenv("provenance_method", "none")
        if fetch_new_documents and documents and provenance_method in ["rerank", "llm", "similarity"]:
            yield ("step", f"Computing provenance scores ({provenance_method})...")
            provenance_scores = self.compute_provenance_scores(prompt, documents, response)

        # Add the response to the history
        new_history.append({"role": "assistant", "content": response})

        # Yield final metadata
        yield ("done", {
            "reply": response,
            "history": new_history,
            "documents": documents,
            "rewritten": rewritten,
            "fetched_new_documents": fetch_new_documents,
            "provenance_scores": provenance_scores,
        })
